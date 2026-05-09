from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def add_title_slide(prs, title, subtitle, stats=None):
    """Slide de titre avec accroche"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 45, 74)  # 0F2D4A

    # Titre principal
    left = Inches(0.5)
    top = Inches(1.2)
    width = Inches(9)
    height = Inches(0.8)

    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Sous-titre bleu
    top = Inches(2.1)
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    p.text = subtitle
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = RGBColor(96, 165, 250)  # 60A5FA
    p.alignment = PP_ALIGN.CENTER

    # Tagline
    top = Inches(2.9)
    height = Inches(0.4)
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    p.text = "De l'experience cachee a l'excellence visible"
    p.font.size = Pt(18)
    p.font.italic = True
    p.font.color.rgb = RGBColor(147, 197, 253)  # 93C5FD
    p.alignment = PP_ALIGN.CENTER

    # Stats boxes si fournies
    if stats:
        y_pos = 4.0
        for stat in stats:
            # Box
            box = slide.shapes.add_shape(
                1, left, Inches(y_pos), Inches(2.8), Inches(1.7)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(22, 43, 66)  # 162B42
            box.line.fill.background()

            # Valeur
            tb = slide.shapes.add_textbox(left, Inches(y_pos + 0.2), Inches(2.8), Inches(0.8))
            tf = tb.text_frame
            p.text = stat["value"]
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = RGBColor(96, 165, 250)
            p.alignment = PP_ALIGN.CENTER

            # Label
            tb = slide.shapes.add_textbox(left, Inches(y_pos + 0.95), Inches(2.8), Inches(0.5))
            tf = tb.text_frame
            p.text = stat["label"]
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(147, 197, 253)
            p.alignment = PP_ALIGN.CENTER

            left += Inches(3.1)
            y_pos = 4.0

    # Footer
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.3))
    tf = tb.text_frame
    p.text = "Formation KAIZEN - Direction Maintenance OCP - 0 Dh d'investissement"
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(75, 107, 138)
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_intro_slide_2(prs):
    """Historique Lean en bref et schematise"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 45, 74)

    # Titre
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    tf = tb.text_frame
    p.text = "INTRODUCTION - HISTORIQUE LEAN"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Sous-titre
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(0.4))
    tf = tb.text_frame
    p.text = "50 ans d'excellence industrielle - une methode eprouvee"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(147, 197, 253)

    # Timeline verticale
    events = [
        ("1950", "Toyota", "TPS - Toyota Production System\nTaiichi Ohno - Le systeme de production", "Toyotisme"),
        ("1970", "Japon", "KAIZEN - Amelioration Continue\nMasaaki Imai - Petites etapes, tous les jours", "Kaizen"),
        ("1986", "USA", "Lean Thinking - Womack & Jones\n5 principes - Eliminer les gaspillages", "Lean"),
        ("1990", "Monde", "Six Sigma - Motorola\nDMAIC - Reduire les variations", "Six Sigma"),
        ("2000", "Industrie", "Industry 4.0 - Lean 4.0\nDigitalisation + Lean + Data", "Lean 4.0"),
    ]

    y_start = 1.8
    for i, (year, place, desc, tag) in enumerate(events):
        y = y_start + i * 0.95

        # Box evenement
        box = slide.shapes.add_shape(
            1, Inches(0.8), Inches(y), Inches(8.4), Inches(0.85)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(22, 43, 66)
        box.line.color.rgb = RGBColor(30, 77, 140)

        # Annee
        tb = slide.shapes.add_textbox(Inches(1.0), Inches(y + 0.1), Inches(0.8), Inches(0.6))
        tf = tb.text_frame
        p.text = year
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = RGBColor(96, 165, 250)

        # Lieu
        tb = slide.shapes.add_textbox(Inches(2.0), Inches(y + 0.15), Inches(1.0), Inches(0.5))
        tf = tb.text_frame
        p.text = place
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        # Description
        tb = slide.shapes.add_textbox(Inches(3.2), Inches(y + 0.1), Inches(4.5), Inches(0.65))
        tf = tb.text_frame
        tf.word_wrap = True
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(200, 200, 200)

        # Tag
        tb = slide.shapes.add_textbox(Inches(7.8), Inches(y + 0.2), Inches(1.2), Inches(0.45))
        tf = tb.text_frame
        p.text = tag
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = RGBColor(96, 165, 250)

    # Message cle
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.5))
    tf = tb.text_frame
    p.text = "🎯 UN HÉRITAGE QUI FONCTIONNE : Les meilleures entreprises industrielles l'appliquent depuis 50 ans"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(96, 165, 250)
    p.alignment = PP_ALIGN.CENTER

    # Slide number
    tb = slide.shapes.add_textbox(Inches(8.5), Inches(7.0), Inches(0.5), Inches(0.3))
    tf = tb.text_frame
    p.text = "2/20"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(147, 197, 253)

def add_part1_slide1(prs):
    """KAIZEN vs KAIKAKU"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 45, 74)

    # Part header
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    tf = tb.text_frame
    p.text = "PARTIE 1 - LEAN THINKING"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Titre
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(0.5))
    tf = tb.text_frame
    p.text = "KAIZEN vs KAIKAKU - Deux voies vers l'amelioration"
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(147, 197, 253)

    # KAIZEN box
    box = slide.shapes.add_shape(
        1, Inches(0.6), Inches(1.8), Inches(4.0), Inches(3.5)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(22, 43, 66)
    box.line.color.rgb = RGBColor(96, 165, 250)
    box.line.width = Pt(3)

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(3.6), Inches(0.5))
    tf = tb.text_frame
    p.text = "改善 KAIZEN"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(96, 165, 250)

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(3.6), Inches(2.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p.text = "Amelioration CONTINUE\n\n• Petites etapes\n• Quotidien\n• Toute l'equipe\n• Sans investissement\n• Effet cumulatif\n\n1% × 365 jours = 37× mieux"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p1 = tf.add_paragraph()
    p1.text = "• Petites etapes"
    p1.level = 0
    p2 = tf.add_paragraph()
    p2.text = "• Quotidien"
    p3 = tf.add_paragraph()
    p3.text = "• Toute l'equipe"
    p4 = tf.add_paragraph()
    p4.text = "• Sans investissement"
    p5 = tf.add_paragraph()
    p5.text = "• Effet cumulatif"
    p6 = tf.add_paragraph()
    p6.text = ""
    p7 = tf.add_paragraph()
    p7.text = "1% x 365 jours = 37x mieux"
    p7.font.color.rgb = RGBColor(96, 165, 250)
    p7.font.bold = True

    # KAIKAKU box
    box = slide.shapes.add_shape(
        1, Inches(5.0), Inches(1.8), Inches(4.0), Inches(3.5)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(45, 20, 20)
    box.line.color.rgb = RGBColor(251, 146, 60)
    box.line.width = Pt(3)

    tb = slide.shapes.add_textbox(Inches(5.2), Inches(2.0), Inches(3.6), Inches(0.5))
    tf = tb.text_frame
    p.text = "改革 KAIKAKU"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(251, 146, 60)

    tb = slide.shapes.add_textbox(Inches(5.2), Inches(2.6), Inches(3.6), Inches(2.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p.text = "Amelioration RADICALE"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p1 = tf.add_paragraph()
    p1.text = ""
    p2 = tf.add_paragraph()
    p2.text = "• Grand changement"
    p3 = tf.add_paragraph()
    p3.text = "• Ponctuel"
    p4 = tf.add_paragraph()
    p4.text = "• Experts/Consultants"
    p5 = tf.add_paragraph()
    p5.text = "• Investissement lourd"
    p6 = tf.add_paragraph()
    p6.text = "• Effet immediat"
    p7 = tf.add_paragraph()
    p7.text = ""
    p8 = tf.add_paragraph()
    p8.text = "Nouvel equipement = 1M Dh"
    p8.font.color.rgb = RGBColor(251, 146, 60)
    p8.font.bold = True

    # VS
    tb = slide.shapes.add_textbox(Inches(4.5), Inches(3.4), Inches(0.5), Inches(0.4))
    tf = tb.text_frame
    p.text = "VS"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Message cle
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(5.6), Inches(9), Inches(0.4))
    tf = tb.text_frame
    p.text = "🎯 NOTRE APPROCHE : KAIZEN - Parce que les petites victoires construisent les grands resultats"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(96, 165, 250)
    p.alignment = PP_ALIGN.CENTER

    # Exemple concret
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(6.2), Inches(8.4), Inches(0.6))
    tf = tb.text_frame
    p.text = "Exemple OCP : Chariot outillage mobile (-15 min/intervention) vs Remplacement concasseur complet (6 mois + 2M Dh)"
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER

    # Slide number
    tb = slide.shapes.add_textbox(Inches(8.5), Inches(7.0), Inches(0.5), Inches(0.3))
    tf = tb.text_frame
    p.text = "3/20"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(147, 197, 253)

def add_part1_slide2(prs):
    """VA / NVA / NvS"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 45, 74)

    # Part header
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.5))
    tf = tb.text_frame
    p.text = "PARTIE 1 - VALEUR AJOUTÉE"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Titre
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(9), Inches(0.4))
    tf = tb.text_frame
    p.text = "Ce que le client paie - et ce qu'il ne paie pas"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(147, 197, 253)

    # Diagramme circulaire simplifie (barres horizontales)
    y = 1.6

    # VA - Vert
    box = slide.shapes.add_shape(
        1, Inches(0.8), Inches(y), Inches(2.5), Inches(0.7)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(34, 197, 94)  # Vert
    box.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(3.4), Inches(y + 0.1), Inches(1.5), Inches(0.5))
    tf = tb.text_frame
    p.text = "VA - 30%"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(34, 197, 94)

    tb = slide.shapes.add_textbox(Inches(5.0), Inches(y + 0.1), Inches(4.0), Inches(0.5))
    tf = tb.text_frame
    p.text = "Valeur Ajoutee - Ce que le client PAIE"
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(255, 255, 255)

    y += 1.0

    # NVA necessaire - Orange
    box = slide.shapes.add_shape(
        1, Inches(0.8), Inches(y), Inches(2.5), Inches(0.7)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(251, 146, 60)  # Orange
    box.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(3.4), Inches(y + 0.1), Inches(1.5), Inches(0.5))
    tf = tb.text_frame
    p.text = "NVA - 25%"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(251, 146, 60)

    tb = slide.shapes.add_textbox(Inches(5.0), Inches(y + 0.1), Inches(4.0), Inches(0.5))
    tf = tb.text_frame
    p.text = "Non-Valeur Ajoutee NÉCESSAIRE - Obligatoire (securite, admin)"
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(255, 255, 255)

    y += 1.0

    # Gaspillage - Rouge
    box = slide.shapes.add_shape(
        1, Inches(0.8), Inches(y), Inches(2.5), Inches(0.7)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(239, 68, 68)  # Rouge
    box.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(3.4), Inches(y + 0.1), Inches(1.8), Inches(0.5))
    tf = tb.text_frame
    p.text = "GASPILLAGE - 45%"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(239, 68, 68)

    tb = slide.shapes.add_textbox(Inches(5.0), Inches(y + 0.1), Inches(4.0), Inches(0.5))
    tf = tb.text_frame
    p.text = "NVA SUPPRIMABLE - Pur gaspillage a ÉLIMINER"
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Exemple concret
    y = 4.8
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(8.4), Inches(1.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p.text = "INTERVENTION MAINTENANCE - Exemple Concasseur (8h)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(96, 165, 250)

    p1 = tf.add_paragraph()
    p1.text = "OK VA (2.5h) : Diagnostic + Reparation + Test - Le client paie pour ca"
    p1.font.size = Pt(12)
    p1.font.color.rgb = RGBColor(34, 197, 94)

    p2 = tf.add_paragraph()
    p2.text = "O NVA Necessaire (1.5h) : Consignation + EPI + Permis feu - Obligatoire"
    p2.font.size = Pt(12)
    p2.font.color.rgb = RGBColor(251, 146, 60)

    p3 = tf.add_paragraph()
    p3.text = "X Gaspillage (4h) : Chercher piece + Attendre + Papier inutile - A SUPPRIMER"
    p3.font.size = Pt(12)
    p3.font.color.rgb = RGBColor(239, 68, 68)

    # Message impact
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.4))
    tf = tb.text_frame
    p.text = "🎯 RÉDUIRE LE GASPILLAGE DE 45% → 20% = +2h/personne/jour = 25% DE CAPACITÉ SUPPLÉMENTAIRE"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(96, 165, 250)
    p.alignment = PP_ALIGN.CENTER

    # Slide number
    tb = slide.shapes.add_textbox(Inches(8.5), Inches(7.0), Inches(0.5), Inches(0.3))
    tf = tb.text_frame
    p.text = "4/20"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(147, 197, 253)

def add_part1_slide3(prs):
    """Les 7 MUDA"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 45, 74)

    # Part header
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.5))
    tf = tb.text_frame
    p.text = "PARTIE 1 - LES 7 MUDA"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Sous-titre
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(9), Inches(0.4))
    tf = tb.text_frame
    p.text = "MUDA = 無駄 (Gaspillage) - Les 7 gaspillages a eliminer en maintenance"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(147, 197, 253)

    # 7 MUDA en grille
    mudas = [
        ("1. SURPRODUCTION", "Rapports non lus\nStocks pieces excessifs\nDouble saisie", "📄"),
        ("2. ATTENTE", "Technicien attend piece\nAttente autorisation\nAttente equipement", "⏳"),
        ("3. TRANSPORT", "Pieces stock loin\nAllers-retours magasin\nDeplacements inutiles", "🚛"),
        ("4. TRAITEMENT EXCESSIF", "Rapports trop detailles\n3 validations\nReunions sans fin", "📝"),
        ("5. STOCKS", "Pieces qui dorment\nObsolescence\nSur-stockage", "📦"),
        ("6. MOUVEMENTS", "Chercher outillage\nMauvaise ergonomie\nAllers-retours", "🏃"),
        ("7. DÉFAUTS", "Reprise - Retouche\nMauvais diagnostic\nRe-reparation", "❌"),
    ]

    y = 1.5
    x = 0.6
    for i, (title, desc, icon) in enumerate(mudas):
        row = i // 2
        col = i % 2

        box_x = Inches(0.6 + col * 4.6)
        box_y = Inches(1.5 + row * 0.85)

        box = slide.shapes.add_shape(1, box_x, box_y, Inches(4.3), Inches(0.75))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(22, 43, 66)
        box.line.color.rgb = RGBColor(30, 77, 140)

        tb = slide.shapes.add_textbox(box_x + Inches(0.1), box_y + Inches(0.05), Inches(4.1), Inches(0.65))
        tf = tb.text_frame
        tf.word_wrap = True
        p.text = f"{icon} {title}\n{desc}"
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(255, 255, 255)

    # 8eme MUDA - Talent ignore
    box = slide.shapes.add_shape(
        1, Inches(0.6), Inches(4.6), Inches(8.8), Inches(0.5)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(96, 165, 50)
    box.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(4.7), Inches(8.4), Inches(0.3))
    tf = tb.text_frame
    p.text = "🧠 8eme MUDA : TALENT IGNORÉ - Ne pas ecouter ceux qui font le travail = gaspiller leur expertise"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Message cle
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(5.4), Inches(9), Inches(0.4))
    tf = tb.text_frame
    p.text = "🎯 CHAQUE MUDA IDENTIFIÉ = DU TEMPS GAGNÉ = DE LA VALEUR CRÉÉE"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(96, 165, 250)
    p.alignment = PP_ALIGN.CENTER

    # Exemple OCP
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(6.0), Inches(8.4), Inches(0.7))
    tf = tb.text_frame
    tf.word_wrap = True
    p.text = "EXEMPLE OCP - 3 MUDA identifies = 75h gagnees/an"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(96, 165, 250)

    p1 = tf.add_paragraph()
    p1.text = "Attente piece (-30h) + Transport pieces (-25h) + Reprise (-20h) = 75h/an sur 1 concasseur"
    p1.font.size = Pt(11)
    p1.font.color.rgb = RGBColor(200, 200, 200)

    # Slide number
    tb = slide.shapes.add_textbox(Inches(8.5), Inches(7.0), Inches(0.5), Inches(0.3))
    tf = tb.text_frame
    p.text = "5/20"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(147, 197, 253)

def add_part2_slides(prs):
    """Partie 2 - Demarches (PDCA, DMAIC, 8D)"""

    # Slide 6 - PDCA
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 45, 74)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.5))
    tf = tb.text_frame
    p.text = "PARTIE 2 - DÉMARCHES"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(9), Inches(0.4))
    tf = tb.text_frame
    p.text = "PDCA - La roue de l'amelioration continue"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(147, 197, 253)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(0.3))
    tf = tb.text_frame
    p.text = "W. Edwards Deming - 1950 - Cycle d'amelioration quotidienne"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(200, 200, 200)

    # 4 quadrants PDCA
    pdca_steps = [
        ("P", "PLAN", "Definir l'objectif\nAnalyser les causes\nPlanifier l'action", RGBColor(59, 130, 246)),
        ("D", "DO", "Tester la solution\nPerimetre limite\nMesurer en continu", RGBColor(34, 197, 94)),
        ("C", "CHECK", "Mesurer les resultats\nAnalyser les ecarts\nValider l'efficacite", RGBColor(251, 146, 60)),
        ("A", "ACT", "Standardiser\nCapitaliser\nEtendre aux autres equipes", RGBColor(168, 85, 247)),
    ]

    positions = [
        (0.6, 2.0), (5.0, 2.0), (0.6, 4.2), (5.0, 4.2)
    ]

    for i, (letter, title, desc, color) in enumerate(pdca_steps):
        x, y = positions[i]

        box = slide.shapes.add_shape(
            1, Inches(x), Inches(y), Inches(4.2), Inches(1.9)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        tb = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.15), Inches(0.5), Inches(0.4))
        tf = tb.text_frame
        p.text = letter
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        tb = slide.shapes.add_textbox(Inches(x + 0.8), Inches(y + 0.2), Inches(3.2), Inches(0.4))
        tf = tb.text_frame
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        tb = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.65), Inches(3.8), Inches(1.1))
        tf = tb.text_frame
        tf.word_wrap = True
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(255, 255, 255)

    # Exemple
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(6.4), Inches(8.8), Inches(0.5))
    tf = tb.text_frame
    p.text = "🎯 EXEMPLE : Arret concasseur 4h/semaine → P: Objectif 1h + Causes: capteur HS + delai pieces → D: Test nouveau capteur → C: Resultat 0.8h → A: Standard sur tous concasseurs"
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(200, 200, 200)

    tb = slide.shapes.add_textbox(Inches(8.5), Inches(7.0), Inches(0.5), Inches(0.3))
    tf = tb.text_frame
    p.text = "6/20"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(147, 197, 253)

    # Slide 7 - DMAIC
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 45, 74)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.5))
    tf = tb.text_frame
    p.text = "PARTIE 2 - DÉMARCHES"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(9), Inches(0.4))
    tf = tb.text_frame
    p.text = "DMAIC - Pour problemes complexes"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(147, 197, 253)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(0.3))
    tf = tb.text_frame
    p.text = "Six Sigma - Motorola 1986 - Projets complexes 1-3 mois"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(200, 200, 200)

    dmaic_steps = [
        ("D", "DEFINIR", "Charte projet\nQQOQCCP\nScope + Objectifs", RGBColor(59, 130, 246)),
        ("M", "MESURER", "Collecter donnees\nBaseline\nIndicateurs", RGBColor(34, 197, 94)),
        ("A", "ANALYSER", "5 Pourquoi\nIshikawa\nDonnees statistiques", RGBColor(251, 146, 60)),
        ("I", "IMPROVE", "Generer solutions\nPilotes\nTest", RGBColor(168, 85, 247)),
        ("C", "CONTROL", "Standardiser\nSuivi KPI\nDocumentation", RGBColor(236, 72, 153)),
    ]

    y = 2.0
    for letter, title, desc, color in dmaic_steps:
        box = slide.shapes.add_shape(
            1, Inches(0.6), Inches(y), Inches(8.8), Inches(0.75)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        tb = slide.shapes.add_textbox(Inches(0.9), Inches(y + 0.1), Inches(0.5), Inches(0.5))
        tf = tb.text_frame
        p.text = letter
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        tb = slide.shapes.add_textbox(Inches(1.5), Inches(y + 0.1), Inches(1.5), Inches(0.5))
        tf = tb.text_frame
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        tb = slide.shapes.add_textbox(Inches(3.2), Inches(y + 0.1), Inches(5.8), Inches(0.55))
        tf = tb.text_frame
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(255, 255, 255)

        y += 0.9

    tb = slide.shapes.add_textbox(Inches(0.6), Inches(6.7), Inches(8.8), Inches(0.4))
    tf = tb.text_frame
    p.text = "UTILISATION : Probleme complexe avec donnees disponibles - Impact financier eleve - Plusieurs equipe concernees"
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(200, 200, 200)

    tb = slide.shapes.add_textbox(Inches(8.5), Inches(7.0), Inches(0.5), Inches(0.3))
    tf = tb.text_frame
    p.text = "7/20"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(147, 197, 253)

    # Slide 8 - 8D
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 45, 74)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.5))
    tf = tb.text_frame
    p.text = "PARTIE 2 - DÉMARCHES"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(9), Inches(0.4))
    tf = tb.text_frame
    p.text = "8D - Resolution d'incidents graves"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(147, 197, 253)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(0.3))
    tf = tb.text_frame
    p.text = "Ford Motor Company - 1987 - Incident grave ou critique - 3 jours a 2 semaines"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(200, 200, 200)

    # 8D en grille 2x4
    eight_d = [
        ("D1", "Équipe", "Multidisciplinaire\nMeca + Élec + Prod"),
        ("D2", "Definir", "Description precise\nQQOQCCP complet"),
        ("D3", "Action immediate", "Stopper l'hemragie\nMesure conservatoire"),
        ("D4", "Causes racines", "5 Pourquoi + Ishikawa\nPreuves"),
        ("D5", "Solutions", "Impact vs Faisabilite\nPermanent"),
        ("D6", "Implementer", "Plan + Responsables\nDelais"),
        ("D7", "Preventives", "Empecher recurrence\nAutres equipements"),
        ("D8", "Capitaliser", "Partager lecons\nBase de connaissances"),
    ]

    y = 1.9
    for i in range(0, 8, 2):
        for j in range(2):
            idx = i + j
            code, title, desc = eight_d[idx]
            x = 0.6 if j == 0 else 4.8

            box = slide.shapes.add_shape(
                1, Inches(x), Inches(y), Inches(4.3), Inches(0.7)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(220, 38, 38) if idx < 3 else RGBColor(22, 43, 66)
            box.line.color.rgb = RGBColor(255, 255, 255)

            tb = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.05), Inches(0.6), Inches(0.6))
            tf = tb.text_frame
            p.text = code
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)

            tb = slide.shapes.add_textbox(Inches(x + 0.8), Inches(y + 0.08), Inches(1.0), Inches(0.25))
            tf = tb.text_frame
            p.text = title
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)

            tb = slide.shapes.add_textbox(Inches(x + 0.8), Inches(y + 0.35), Inches(3.2), Inches(0.3))
            tf = tb.text_frame
            tf.word_wrap = True
            p.text = desc
            p.font.size = Pt(9)
            p.font.color.rgb = RGBColor(200, 200, 200)

        y += 0.85

    tb = slide.shapes.add_textbox(Inches(0.6), Inches(5.0), Inches(8.8), Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p.text = "🎯 EXEMPLE OCP - Courroie transporteur rompue 3x en 1 mois"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RGBColor(251, 146, 60)

    p1 = tf.add_paragraph()


    p1.text = "D1: Équipe meca+elec+prod | D2: Courroie BW2 rompue | D3: Installation temporaire | D4: Tension excessive + qualite courroie | D5-D7: Nouveau type + plan maintenance | D8: Standard tous transporteurs"
    p1.font.size = Pt(10)

    tb = slide.shapes.add_textbox(Inches(8.5), Inches(7.0), Inches(0.5), Inches(0.3))
    tf = tb.text_frame
    p.text = "8/20"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(147, 197, 253)

def add_part3_slides(prs):
    """Partie 3 - Outils terrain"""

    # Slide 9 - QQOQCCP
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 45, 74)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.5))
    tf = tb.text_frame
    p.text = "PARTIE 3 - OUTILS TERRAIN"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(9), Inches(0.4))
    tf = tb.text_frame
    p.text = "QQOQCCP - Comprendre avant d'agir"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(147, 197, 253)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(0.3))
    tf = tb.text_frame
    p.text = "Methode journalistique - 30 minutes - Pour bien definir un probleme"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(200, 200, 200)

    # QQOQCCP grid
    qqoqccp = [
        ("Q", "Qui ?", "Équipe A - Shift matin - Technicien: M. Benali"),
        ("Q", "Quoi ?", "Vibration anormale broyeur - Bruit inhabituel"),
        ("O", "Où ?", "Broyeur BK3 - Palier 2 - Cote moteur"),
        ("Q", "Quand ?", "Demarrage production - 6h du matin - Depuis 3 jours"),
        ("C", "Comment ?", "Vibration 12mm/s - Alarme rouge - Temperature +15°C"),
        ("C", "Combien ?", "3x/semaine - -2h/arret - -500t production"),
        ("P", "Pourquoi ?", "Perte production 500t/an - Risque casseur - Cout 20k Dh"),
    ]

    y = 1.8
    for letter, question, answer in qqoqccp:
        # Question box
        box = slide.shapes.add_shape(
            1, Inches(0.6), Inches(y), Inches(1.0), Inches(0.55)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(96, 165, 250)
        box.line.fill.background()

        tb = slide.shapes.add_textbox(Inches(0.7), Inches(y + 0.1), Inches(0.8), Inches(0.35))
        tf = tb.text_frame
        p.text = f"{letter}\n{question}"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        # Answer box
        box = slide.shapes.add_textbox(Inches(1.7), Inches(y), Inches(7.5), Inches(0.55))
        tf = box.text_frame
        tf.word_wrap = True
        p.text = answer
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(200, 200, 200)

        y += 0.65

    tb = slide.shapes.add_textbox(Inches(0.6), Inches(6.5), Inches(8.8), Inches(0.4))
    tf = tb.text_frame
    p.text = "🎯 RÉSULTAT : Probleme bien defini = Solution adaptee = Plus de recurrence"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RGBColor(96, 165, 250)

    tb = slide.shapes.add_textbox(Inches(8.5), Inches(7.0), Inches(0.5), Inches(0.3))
    tf = tb.text_frame
    p.text = "9/20"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(147, 197, 253)

    # Slide 10 - 5 Pourquoi
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 45, 74)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.5))
    tf = tb.text_frame
    p.text = "PARTIE 3 - OUTILS TERRAIN"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(9), Inches(0.4))
    tf = tb.text_frame
    p.text = "5 POURQUOI - Trouver la cause racine"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(147, 197, 253)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(0.3))
    tf = tb.text_frame
    p.text = "Taiichi Ohno - Toyota 1950 - 15 minutes - Permet de descendre a la racine du probleme"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(200, 200, 200)

    # Arbre des pourquoi
    problem = "Moteur convoyeur L4 en panne"
    effect = "Arret ligne 4h - Perte 40t - Cout 8k Dh"

    tb = slide.shapes.add_textbox(Inches(0.6), Inches(1.7), Inches(8.8), Inches(0.4))
    tf = tb.text_frame
    p.text = f"PROBLÈME : {problem}"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(239, 68, 68)

    tb = slide.shapes.add_textbox(Inches(5.5), Inches(1.7), Inches(3.9), Inches(0.4))
    tf = tb.text_frame
    p.text = f"→ {effect}"
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(251, 146, 60)

    y = 2.3
    pourquoi_steps = [
        ("1", "Pourquoi panne ?", "Surchauffe moteur"),
        ("2", "Pourquoi surchauffe ?", "Ventilation encrassee"),
        ("3", "Pourquoi encrassee ?", "Filtre absent depuis 3 mois"),
        ("4", "Pourquoi absent ?", "Pas de plan maintenance preventive"),
        ("5", "CAUSE RACINE", "Systeme de gestion MP inexistant"),
    ]

    for num, question, answer in pourquoi_steps:
        # Arrow/step
        tb = slide.shapes.add_textbox(Inches(1.5), Inches(y), Inches(0.5), Inches(0.4))
        tf = tb.text_frame
        p.text = "↓"
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(147, 197, 253)

        tb = slide.shapes.add_textbox(Inches(0.6), Inches(y + 0.1), Inches(3.5), Inches(0.4))
        tf = tb.text_frame
        p.text = f"{num}. {question}"
        p.font.size = Pt(13)
        p.font.color.rgb = RGBColor(255, 255, 255)

        tb = slide.shapes.add_textbox(Inches(4.2), Inches(y + 0.1), Inches(5.0), Inches(0.4))
        tf = tb.text_frame
        p.text = f"→ {answer}"
        p.font.size = Pt(12)

        if num == "5":
            p.font.bold = True
            p.font.color.rgb = RGBColor(34, 197, 94)
        else:
            p.font.color.rgb = RGBColor(200, 200, 200)

        y += 0.5

    # Solution box
    box = slide.shapes.add_shape(
        1, Inches(0.6), Inches(5.0), Inches(8.8), Inches(1.2)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(34, 197, 94)
    box.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(5.2), Inches(8.4), Inches(0.9))
    tf = tb.text_frame
    tf.word_wrap = True
    p.text = "SOLUTION + RESULTAT"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    p1 = tf.add_paragraph()
    p1.text = "SOLUTION : Implementer plan MP - GEMBA mensuel - Check-liste ventilation"
    p1.font.size = Pt(11)
    p1.font.color.rgb = RGBColor(255, 255, 255)

    p2 = tf.add_paragraph()
    p2.text = "RESULTAT : 0 panne ventilation en 6 mois - Gain 24h/an - Economie 12k Dh"
    p2.font.size = Pt(11)
    p2.font.color.rgb = RGBColor(200, 255, 200)

    tb = slide.shapes.add_textbox(Inches(8.5), Inches(7.0), Inches(0.5), Inches(0.3))
    tf = tb.text_frame
    p.text = "10/20"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(147, 197, 253)

    # Slide 11 - SMED
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 45, 74)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.5))
    tf = tb.text_frame
    p.text = "PARTIE 3 - OUTILS TERRAIN"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(9), Inches(0.4))
    tf = tb.text_frame
    p.text = "SMED - Reduire temps de changement"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(147, 197, 253)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(0.3))
    tf = tb.text_frame
    p.text = "Single Minute Exchange of Die - Shigeo Shingo 1950 - Projet 1 mois - -50% temps changement"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(200, 200, 200)

    # INTERNE vs EXTERNE
    # INTERNE box
    box = slide.shapes.add_shape(
        1, Inches(0.6), Inches(1.8), Inches(4.0), Inches(2.0)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(239, 68, 68)
    box.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(3.6), Inches(0.4))
    tf = tb.text_frame
    p.text = "INTERNE (machine arretee)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(3.6), Inches(1.1))
    tf = tb.text_frame
    tf.word_wrap = True
    p.text = "• Demontage\n• Nettoyage\n• Installation\n• Alignement\n• Test"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(255, 255, 255)

    # EXTERNE box
    box = slide.shapes.add_shape(
        1, Inches(5.0), Inches(1.8), Inches(4.0), Inches(2.0)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(34, 197, 94)
    box.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(5.2), Inches(2.0), Inches(3.6), Inches(0.4))
    tf = tb.text_frame
    p.text = "EXTERNE (machine en marche)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    tb = slide.shapes.add_textbox(Inches(5.2), Inches(2.5), Inches(3.6), Inches(1.1))
    tf = tb.text_frame
    tf.word_wrap = True
    p.text = "• Preparer outils\n• Prechauffer\n• À portee main\n• Check-list prete"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Arrow
    tb = slide.shapes.add_textbox(Inches(4.4), Inches(2.6), Inches(0.4), Inches(0.4))
    tf = tb.text_frame
    p.text = "→"
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(96, 165, 250)
    p.alignment = PP_ALIGN.CENTER

    tb = slide.shapes.add_textbox(Inches(4.2), Inches(2.9), Inches(0.8), Inches(0.3))
    tf = tb.text_frame
    p.text = "Transformer"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(96, 165, 250)
    p.alignment = PP_ALIGN.CENTER

    # Exemple
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(4.1), Inches(8.8), Inches(1.3))
    tf = tb.text_frame
    tf.word_wrap = True
    p.text = "🎯 EXEMPLE OCP - Changement males concasseur"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(96, 165, 250)

    p1 = tf.add_paragraph()


    p1.text = "AVANT : 4h 30min (INTERNE)\nAPRÈS : 2h 15min (SMED)\nGAIN : -50% - 2h/fois - 150h/an sur 3 concasseurs"
    p1.font.size = Pt(12)

    p2 = tf.add_paragraph()
    p2.font.size = Pt(11)

    # Message
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(5.6), Inches(8.8), Inches(0.4))
    tf = tb.text_frame
    p.text = "💡 PRINCIPE : Transformer maximum de taches INTERNES en EXTERNES = Reduire arret machine"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RGBColor(34, 197, 94)

    tb = slide.shapes.add_textbox(Inches(8.5), Inches(7.0), Inches(0.5), Inches(0.3))
    tf = tb.text_frame
    p.text = "11/20"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(147, 197, 253)

    # Slide 12 - Brainstorming
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 45, 74)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.5))
    tf = tb.text_frame
    p.text = "PARTIE 3 - OUTILS TERRAIN"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(9), Inches(0.4))
    tf = tb.text_frame
    p.text = "BRAINSTORMING - Generer des solutions"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(147, 197, 253)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(0.3))
    tf = tb.text_frame
    p.text = "4 formats adaptes terrain maintenance"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(200, 200, 200)

    formats = [
        ("CLASSIQUE", "6-10 pers", "20-40 min", "Reunion equipe - Animateur - Tableau blanc", RGBColor(59, 130, 246)),
        ("STAND-UP", "5-8 pers", "10-15 min", "Debout - Rapide - 1 probleme unique", RGBColor(34, 197, 94)),
        ("BRAIN WRITING", "6 pers", "30 min", "Chacun ecrit 3 idees - Passe feuille", RGBColor(168, 85, 247)),
        ("SOLO", "1 pers", "15 min", "Seul - Note TOUS les gaspillages vus", RGBColor(251, 146, 60)),
    ]

    y = 1.8
    for title, people, time, desc, color in formats:
        box = slide.shapes.add_shape(
            1, Inches(0.6), Inches(y), Inches(8.8), Inches(0.85)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        tb = slide.shapes.add_textbox(Inches(0.8), Inches(y + 0.1), Inches(2.5), Inches(0.3))
        tf = tb.text_frame
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        tb = slide.shapes.add_textbox(Inches(3.4), Inches(y + 0.1), Inches(1.2), Inches(0.3))
        tf = tb.text_frame
        p.text = people
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(255, 255, 255)

        tb = slide.shapes.add_textbox(Inches(4.7), Inches(y + 0.1), Inches(1.0), Inches(0.3))
        tf = tb.text_frame
        p.text = time
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(255, 255, 255)

        tb = slide.shapes.add_textbox(Inches(5.8), Inches(y + 0.1), Inches(3.4), Inches(0.65))
        tf = tb.text_frame
        tf.word_wrap = True
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(255, 255, 255)

        y += 1.0

    tb = slide.shapes.add_textbox(Inches(0.6), Inches(5.9), Inches(8.8), Inches(0.4))
    tf = tb.text_frame
    p.text = "🎯 EXEMPLE SOLO : Technicien note 12 gaspillages pendant son shift → Best : Detection precoce 3 pannes evitees"
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(200, 200, 200)

    tb = slide.shapes.add_textbox(Inches(8.5), Inches(7.0), Inches(0.5), Inches(0.3))
    tf = tb.text_frame
    p.text = "12/20"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(147, 197, 253)

    # Slide 13 - 5G & Andon
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 45, 74)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.5))
    tf = tb.text_frame
    p.text = "PARTIE 3 - OUTILS 2ÈME GÉNÉRATION"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(9), Inches(0.4))
    tf = tb.text_frame
    p.text = "5G - ANDON - Management visuel"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(147, 197, 253)

    # 5G box
    box = slide.shapes.add_shape(
        1, Inches(0.6), Inches(1.5), Inches(4.0), Inches(2.8)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(22, 43, 66)
    box.line.color.rgb = RGBColor(96, 165, 250)

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(3.6), Inches(0.4))
    tf = tb.text_frame
    p.text = "5G - Les 5 Gemba"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(96, 165, 250)

    five_g = [
        ("1. GEMBA", "Aller sur terrain - Atelier - Ligne de production"),
        ("2. GEMBUTSU", "Observer l'equipement reel - Toucher - Écouter"),
        ("3. GENJITSU", "Donnees reelles - MES - Pas de supposition"),
        ("4. GENRI", "Principes ingenierie - Physique - Mecanique"),
        ("5. GEMBA WALK", "15 min/jour - Noter 3 MUDA"),
    ]

    y = 2.2
    for title, desc in five_g:
        tb = slide.shapes.add_textbox(Inches(0.9), Inches(y), Inches(3.4), Inches(0.3))
        tf = tb.text_frame
        p.text = f"{title} : {desc}"
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(255, 255, 255)

        y += 0.35

    # ANDON box
    box = slide.shapes.add_shape(
        1, Inches(5.0), Inches(1.5), Inches(4.2), Inches(2.8)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(22, 43, 66)
    box.line.color.rgb = RGBColor(96, 165, 250)

    tb = slide.shapes.add_textbox(Inches(5.2), Inches(1.7), Inches(3.8), Inches(0.4))
    tf = tb.text_frame
    p.text = "ANDON - Management Visuel"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(96, 165, 250)

    # ANDON colors
    andon_colors = [
        ("VERT", "Équipement OK - Production normale", RGBColor(34, 197, 94)),
        ("ORANGE", "Vigilance - Planifie maintenance", RGBColor(251, 146, 60)),
        ("ROUGE", "Panne - Action immediate requise", RGBColor(239, 68, 68)),
    ]

    y = 2.2
    for title, desc, color in andon_colors:
        box = slide.shapes.add_shape(
            1, Inches(5.2), Inches(y), Inches(3.8), Inches(0.35)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        tb = slide.shapes.add_textbox(Inches(5.4), Inches(y + 0.05), Inches(3.4), Inches(0.25))
        tf = tb.text_frame
        p.text = f"{title} - {desc}"
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(255, 255, 255)

        y += 0.45

    # Exemple
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(4.5), Inches(8.8), Inches(0.4))
    tf = tb.text_frame
    p.text = "🎯 EXEMPLE OCP - Tableau ANDON salle controle : 20 equipements - Vert/Orange/Rouge visible de loin - Temps reaction panne -40%"
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(200, 200, 200)

    tb = slide.shapes.add_textbox(Inches(8.5), Inches(7.0), Inches(0.5), Inches(0.3))
    tf = tb.text_frame
    p.text = "13/20"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(147, 197, 253)

    # Slide 14 - 5G vs Gensoku
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 45, 74)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.5))
    tf = tb.text_frame
    p.text = "PARTIE 3 - COMPARAISON 5G / GENGSOKU"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(9), Inches(0.4))
    tf = tb.text_frame
    p.text = "Comprendre l'ecart pour ajuster votre pratique"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(147, 197, 253)

    # 5G box
    box = slide.shapes.add_shape(
        1, Inches(0.6), Inches(1.5), Inches(4.0), Inches(3.5)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(22, 43, 66)
    box.line.color.rgb = RGBColor(96, 165, 250)

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(3.6), Inches(0.4))
    tf = tb.text_frame
    p.text = "5G - OBSERVER"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(96, 165, 250)

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(3.6), Inches(2.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p.text = "• Aller sur terrain (Gemba)\n• Voir l'equipement (Gembutsu)\n• Donnees reelles (Genjitsu)\n• Principes (Genri)\n• Walk 15 min/jour"
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(255, 255, 255)

    p1 = tf.add_paragraph()


    p1.text = "But : OBSERVER la realite\nComprendre ce qui se passe\nreellement sur le terrain"
    p1.font.size = Pt(11)
    p1 = tf.add_paragraph()
    p1.font.italic = True

    # GENGSOKU box
    box = slide.shapes.add_shape(
        1, Inches(5.0), Inches(1.5), Inches(4.0), Inches(3.5)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(45, 30, 20)
    box.line.color.rgb = RGBColor(251, 146, 60)

    tb = slide.shapes.add_textbox(Inches(5.2), Inches(1.7), Inches(3.6), Inches(0.4))
    tf = tb.text_frame
    p.text = "GENGSOKU - EXPÉRIMENTER"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(251, 146, 60)

    tb = slide.shapes.add_textbox(Inches(5.2), Inches(2.2), Inches(3.6), Inches(2.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p.text = "• Faire soi-meme\n• Essayer l'operation\n• Ressentir la difficulte\n• Identifier freins reels\n• Duree : 1-2h"
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(255, 255, 255)

    p1 = tf.add_paragraph()


    p1.text = "But : VIVRE le probleme\nExperimenter ce que\nles operateurs subissent"
    p1.font.size = Pt(11)
    p1 = tf.add_paragraph()
    p1.font.italic = True

    # VS arrow
    tb = slide.shapes.add_textbox(Inches(4.4), Inches(3.0), Inches(0.4), Inches(0.4))
    tf = tb.text_frame
    p.text = "+"
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Message
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(5.2), Inches(8.8), Inches(0.4))
    tf = tb.text_frame
    p.text = "🎯 5G = Vue manager - GENGSOKU = Vecu operateur → COMPLÉMENTAIRES pour une vraie comprehension"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RGBColor(96, 165, 250)
    p.alignment = PP_ALIGN.CENTER

    tb = slide.shapes.add_textbox(Inches(0.6), Inches(5.7), Inches(8.8), Inches(0.6))
    tf = tb.text_frame
    p.text = "EXEMPLE : 5G → Je vois le technicien chercher l'outil | GENGSOKU → Je cherche l'outil moi-meme et je perds 15 min → ACTION : Chariot outillage mobile"
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(200, 200, 200)

    tb = slide.shapes.add_textbox(Inches(8.5), Inches(7.0), Inches(0.5), Inches(0.3))
    tf = tb.text_frame
    p.text = "14/20"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(147, 197, 253)

def add_part4_slides(prs):
    """Partie 4 - Gains et Plan d'action"""

    # Slide 15 - Gains quantitatifs
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 45, 74)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.5))
    tf = tb.text_frame
    p.text = "PARTIE 4 - GAINS"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(9), Inches(0.4))
    tf = tb.text_frame
    p.text = "Ce que vous gagnez reellement - Resultats mesures OCP + industrie"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(147, 197, 253)

    # Gains boxes
    gains = [
        ("-40%", "Pannes recurrentes", "En 6 mois avec PDCA + 5P", RGBColor(34, 197, 94)),
        ("+2h/j", "Disponibilite", "Par equipement - Sans investissement", RGBColor(59, 130, 246)),
        ("-50%", "Temps d'intervention", "Avec SMED + Organisation", RGBColor(168, 85, 247)),
        ("0 Dh", "D'investissement", "Juste votre regard - 15 min/jour", RGBColor(251, 146, 60)),
    ]

    positions = [
        (0.6, 1.7), (5.0, 1.7), (0.6, 3.5), (5.0, 3.5)
    ]

    for i, (value, title, desc, color) in enumerate(gains):
        x, y = positions[i]

        box = slide.shapes.add_shape(
            1, Inches(x), Inches(y), Inches(4.2), Inches(1.6)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(22, 43, 66)
        box.line.color.rgb = color
        box.line.width = Pt(2)

        tb = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.2), Inches(3.8), Inches(0.5))
        tf = tb.text_frame
        p.text = value
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER

        tb = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.75), Inches(3.8), Inches(0.35))
        tf = tb.text_frame
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        tb = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 1.1), Inches(3.8), Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(200, 200, 200)
        p.alignment = PP_ALIGN.CENTER

    # Timeline
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(5.4), Inches(8.8), Inches(0.4))
    tf = tb.text_frame
    p.text = "📈 TIMELANE DE RÉSULTATS"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(96, 165, 250)

    timeline = [
        ("1ere SEMAINE", "Identifier MUDA - Gemba Walk"),
        ("1er MOIS", "Premiers resultats visibles"),
        ("3 MOIS", "Équipe autonome"),
        ("6 MOIS", "-40% pannes recurrentes"),
    ]

    y = 5.9
    for period, result in timeline:
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(8.4), Inches(0.25))
        tf = tb.text_frame
        p.text = f"{period} → {result}"
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(255, 255, 255)

        y += 0.3

    tb = slide.shapes.add_textbox(Inches(8.5), Inches(7.0), Inches(0.5), Inches(0.3))
    tf = tb.text_frame
    p.text = "15/20"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(147, 197, 253)

    # Slide 16 - Plan d'action 18-22 mai
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 45, 74)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.5))
    tf = tb.text_frame
    p.text = "PARTIE 4 - PLAN D'ACTION"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(9), Inches(0.4))
    tf = tb.text_frame
    p.text = "Semaine pratique - Du 18 au 22 Mai 2026"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(147, 197, 253)

    # 5 jours
    week_plan = [
        ("LUNDI 18", "5G - Gemba Walk", "15 min terrain\nNoter 3 MUDA observes", "3 MUDA identifies", RGBColor(59, 130, 246)),
        ("MARDI 19", "QQOQCCP + 5 Pourquoi", "30 min sur MUDA critique\nÉquipe en stand-up", "Cause racine trouvee", RGBColor(34, 197, 94)),
        ("MERCREDI 20", "Brainstorming", "10 min stand-up\nÉquipe propose solutions", "Solution choisie", RGBColor(168, 85, 247)),
        ("JEUDI 21", "PDCA - Do", "Tester solution\nMesurer avant/apres", "Resultat mesure", RGBColor(251, 146, 60)),
        ("VENDREDI 22", "Standardisation", "Si succes : Standard\nPartager aux autres equipes", "Standard diffuse", RGBColor(34, 197, 94)),
    ]

    y = 1.5
    for day, tool, action, result, color in week_plan:
        box = slide.shapes.add_shape(
            1, Inches(0.6), Inches(y), Inches(8.8), Inches(1.0)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(22, 43, 66)
        box.line.color.rgb = color

        tb = slide.shapes.add_textbox(Inches(0.8), Inches(y + 0.1), Inches(1.5), Inches(0.3))
        tf = tb.text_frame
        p.text = day
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = color

        tb = slide.shapes.add_textbox(Inches(2.4), Inches(y + 0.1), Inches(2.0), Inches(0.3))
        tf = tb.text_frame
        p.text = tool
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        tb = slide.shapes.add_textbox(Inches(0.8), Inches(y + 0.45), Inches(3.5), Inches(0.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p.text = action
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(200, 200, 200)

        # Result box
        box2 = slide.shapes.add_shape(
            1, Inches(5.5), Inches(y + 0.15), Inches(3.7), Inches(0.7)
        )
        box2.fill.solid()
        box2.fill.fore_color.rgb = color
        box2.line.fill.background()

        tb = slide.shapes.add_textbox(Inches(5.7), Inches(y + 0.25), Inches(3.3), Inches(0.5))
        tf = tb.text_frame
        p.text = f"✓ {result}"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        y += 1.15

    # Footer
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(7.2), Inches(8.8), Inches(0.4))
    tf = tb.text_frame
    p.text = "🎯 OBJECTIF SEMAINE : 1 KAIZEN teste par equipe = 5 ameliorations concretes OCP"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(96, 165, 250)
    p.alignment = PP_ALIGN.CENTER

    tb = slide.shapes.add_textbox(Inches(8.5), Inches(7.5), Inches(0.5), Inches(0.3))
    tf = tb.text_frame
    p.text = "16/20"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(147, 197, 253)

def add_annex_slides(prs):
    """Annexes - Memo et conclusion"""

    # Slide 17 - Memo outils
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 45, 74)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.5))
    tf = tb.text_frame
    p.text = "ANNEXE - MÉMO OUTILS TERRAIN"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(9), Inches(0.3))
    tf = tb.text_frame
    p.text = "À emporter - Guide rapide terrain"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(147, 197, 253)

    # Outils grid
    tools = [
        ("PDCA", "Amelioration quotidienne", "Cycle 1 semaine"),
        ("QQOQCCP", "Definir probleme", "30 min"),
        ("5 Pourquoi", "Cause racine", "15 min"),
        ("SMED", "Reduire temps changement", "Projet 1 mois"),
        ("Brainstorming", "Generer solutions", "10-40 min"),
        ("5G/Andon", "Management visuel", "15 min/jour"),
        ("8D", "Incident grave", "3j-2sem"),
        ("DMAIC", "Projet complexe", "1-3 mois"),
    ]

    y = 1.4
    for i in range(0, 8, 2):
        for j in range(2):
            idx = i + j
            tool, usage, duree = tools[idx]
            x = 0.6 if j == 0 else 4.8

            box = slide.shapes.add_shape(
                1, Inches(x), Inches(y), Inches(4.3), Inches(0.65)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(22, 43, 66)
            box.line.color.rgb = RGBColor(96, 165, 250)

            tb = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.08), Inches(1.0), Inches(0.25))
            tf = tb.text_frame
            p.text = tool
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = RGBColor(96, 165, 250)

            tb = slide.shapes.add_textbox(Inches(x + 1.2), Inches(y + 0.08), Inches(1.5), Inches(0.25))
            tf = tb.text_frame
            p.text = usage
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(255, 255, 255)

            tb = slide.shapes.add_textbox(Inches(x + 2.8), Inches(y + 0.08), Inches(1.3), Inches(0.25))
            tf = tb.text_frame
            p.text = duree
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(200, 200, 200)

        y += 0.75

    # Checklist quotidien
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(5.2), Inches(8.8), Inches(0.3))
    tf = tb.text_frame
    p.text = "✅ CHECKLIST QUOTIDIEN RESPONSABLE"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(96, 165, 250)

    checklist = [
        "□ Gemba Walk 15 min - Noter 3 MUDA",
        "□ Verifier ANDON - Planifier actions ROUGES",
        "□ Briefing equipe - QQOQCCP sur probleme du jour",
        "□ Valider 1 amelioration KAIZEN",
    ]

    y = 5.6
    for item in checklist:
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(8.4), Inches(0.25))
        tf = tb.text_frame
        p.text = item
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(255, 255, 255)
        y += 0.3

    tb = slide.shapes.add_textbox(Inches(8.5), Inches(7.0), Inches(0.5), Inches(0.3))
    tf = tb.text_frame
    p.text = "17/20"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(147, 197, 253)

    # Slide 18 - Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 45, 74)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.5))
    tf = tb.text_frame
    p.text = "CONCLUSION"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Citation
    tb = slide.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(8.0), Inches(0.8))
    tf = tb.text_frame
    p.text = """"Lean Maintenance : Un voyage, pas une destination."""
    p.font.size = Pt(22)
    p.font.italic = True
    p.font.color.rgb = RGBColor(96, 165, 250)
    p.alignment = PP_ALIGN.CENTER

    # 3 principes
    principes = [
        ("Commencez PETIT", "1 MUDA a la fois"),
        ("Commencez MAINTENANT", "Lundi 18 mai"),
        ("Capitalisez TOUJOURS", "Partagez les succes"),
    ]

    y = 2.3
    for principe, desc in principes:
        box = slide.shapes.add_shape(
            1, Inches(1.5), Inches(y), Inches(7.0), Inches(0.6)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(22, 43, 66)
        box.line.color.rgb = RGBColor(96, 165, 250)

        tb = slide.shapes.add_textbox(Inches(1.7), Inches(y + 0.1), Inches(3.0), Inches(0.4))
        tf = tb.text_frame
        p.text = principe
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(96, 165, 250)

        tb = slide.shapes.add_textbox(Inches(4.8), Inches(y + 0.1), Inches(3.2), Inches(0.4))
        tf = tb.text_frame
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(255, 255, 255)

        y += 0.8

    # Message final
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(4.8), Inches(8.8), Inches(0.5))
    tf = tb.text_frame
    p.text = "Votre equipe vous attend. Les outils sont prets. La premiere amelioration commence lundi."
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(34, 197, 94)
    p.alignment = PP_ALIGN.CENTER

    # Contact
    tb = slide.shapes.add_textbox(Inches(2.0), Inches(6.0), Inches(6.0), Inches(0.4))
    tf = tb.text_frame
    p.text = "Direction Maintenance OCP - Formation KAIZEN"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER

    tb = slide.shapes.add_textbox(Inches(8.5), Inches(7.0), Inches(0.5), Inches(0.3))
    tf = tb.text_frame
    p.text = "18/20"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(147, 197, 253)

# Creer la presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# INTRODUCTION (2 slides)
stats = [
    {"value": "-40%", "label": "Pannes recurrentes"},
    {"value": "+2h/j", "label": "Disponibilite equipement"},
    {"value": "0 Dh", "label": "D'investissement"},
]
add_title_slide(prs, "Maintenance Lean &\nAmelioration Continue", "Formation KAIZEN - Direction Maintenance OCP", stats)
add_intro_slide_2(prs)

# PARTIE 1 (3 slides)
add_part1_slide1(prs)  # KAIZEN vs KAIKAKU
add_part1_slide2(prs)  # VA/NVA/NvS
add_part1_slide3(prs)  # 7 MUDA

# PARTIE 2 (3 slides)
add_part2_slides(prs)

# PARTIE 3 (6 slides)
add_part3_slides(prs)

# PARTIE 4 (2 slides)
add_part4_slides(prs)

# ANNEXES (2 slides)
add_annex_slides(prs)

# Sauvegarder
output_path = "C:/Users/HP/Desktop/formation/Formation_Lean_KAIZEN_OCP_V2.pptx"
prs.save(output_path)
print(f"Presentation creee : {output_path}")
