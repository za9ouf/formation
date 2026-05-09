from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def set_paragraph(tf, index, text, font_size=11, bold=False, color=None):
    """Helper pour définir un paragraphe"""
    while len(tf.paragraphs) <= index:
        tf.add_paragraph()
    p = tf.paragraphs[index]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    if color:
        p.font.color.rgb = color
    return p

def add_simple_slide(prs, title, content_lines, subtitle=""):
    """Ajouter un slide simple avec titre et contenu"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 45, 74)

    # Titre
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Sous-titre
    if subtitle:
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(0.3))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(147, 197, 253)

    # Contenu
    y = 1.5
    for line_text in content_lines:
        tb = slide.shapes.add_textbox(Inches(0.6), Inches(y), Inches(8.8), Inches(0.5))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = line_text
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(220, 220, 220)
        y += 0.6

    return slide

# Creer la presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# ===== SLIDE 1 : TITRE =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(15, 45, 74)

tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.8))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "Maintenance Lean &\nAmelioration Continue"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(9), Inches(0.5))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "Formation KAIZEN - Direction Maintenance OCP"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = RGBColor(96, 165, 250)
p.alignment = PP_ALIGN.CENTER

tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(9), Inches(0.4))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "De l'experience cachee a l'excellence visible - 0 Dh d'investissement"
p.font.size = Pt(16)
p.font.italic = True
p.font.color.rgb = RGBColor(147, 197, 253)
p.alignment = PP_ALIGN.CENTER

# Stats
stats = [
    ("-40%", "Pannes recurrentes"),
    ("+2h/j", "Disponibilite equipement"),
    ("0 Dh", "D'investissement"),
]
x = 0.6
for value, label in stats:
    box = slide.shapes.add_shape(1, Inches(x), Inches(4.0), Inches(2.8), Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(22, 43, 66)
    box.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(x), Inches(4.2), Inches(2.8), Inches(0.6))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = value
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(96, 165, 250)
    p.alignment = PP_ALIGN.CENTER

    tb = slide.shapes.add_textbox(Inches(x), Inches(4.8), Inches(2.8), Inches(0.4))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(147, 197, 253)
    p.alignment = PP_ALIGN.CENTER

    x += 3.1

tb = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.3))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "1/18"
p.font.size = Pt(10)
p.font.color.rgb = RGBColor(147, 197, 253)
p.alignment = PP_ALIGN.CENTER

# ===== SLIDE 2 : HISTORIQUE =====
add_simple_slide(prs, "INTRODUCTION - HISTORIQUE LEAN",
    ["1950 - Toyota : TPS (Toyota Production System) par Taiichi Ohno",
     "1970 - Japon : KAIZEN (Amelioration Continue) par Masaaki Imai",
     "1986 - USA : Lean Thinking par Womack & Jones",
     "1990 - Monde : Six Sigma par Motorola",
     "2000 - Industrie : Industry 4.0 - Lean 4.0"],
    "50 ans d'excellence industrielle - Une methode eprouvee")

# ===== SLIDE 3 : KAIZEN VS KAIKAKU =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(15, 45, 74)

tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.5))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "PARTIE 1 - LEAN THINKING"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)

tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(0.4))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "KAIZEN vs KAIKAKU - Deux voies vers l'amelioration"
p.font.size = Pt(18)
p.font.color.rgb = RGBColor(147, 197, 253)

# KAIZEN box
box = slide.shapes.add_shape(1, Inches(0.6), Inches(1.6), Inches(4.0), Inches(3.0))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(22, 43, 66)
box.line.color.rgb = RGBColor(96, 165, 250)

tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(3.6), Inches(2.6))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "KAIZEN - Amelioration Continue"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = RGBColor(96, 165, 250)

p = tf.add_paragraph()
p.text = "Petites etapes - Quotidien - Toute l'equipe"
p.font.size = Pt(12)
p.font.color.rgb = RGBColor(255, 255, 255)

p = tf.add_paragraph()
p.text = "Sans investissement - Effet cumulatif"
p.font.size = Pt(12)
p.font.color.rgb = RGBColor(255, 255, 255)

p = tf.add_paragraph()
p.text = "1% x 365 jours = 37x mieux"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = RGBColor(96, 165, 250)

# KAIKAKU box
box = slide.shapes.add_shape(1, Inches(5.0), Inches(1.6), Inches(4.0), Inches(3.0))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(45, 20, 20)
box.line.color.rgb = RGBColor(251, 146, 60)

tb = slide.shapes.add_textbox(Inches(5.2), Inches(1.8), Inches(3.6), Inches(2.6))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "KAIKAKU - Amelioration Radicale"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = RGBColor(251, 146, 60)

p = tf.add_paragraph()
p.text = "Grand changement - Ponctuel - Experts"
p.font.size = Pt(12)
p.font.color.rgb = RGBColor(255, 255, 255)

p = tf.add_paragraph()
p.text = "Investissement lourd - Effet immediat"
p.font.size = Pt(12)
p.font.color.rgb = RGBColor(255, 255, 255)

p = tf.add_paragraph()
p.text = "Nouvel equipement = 1M Dh"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = RGBColor(251, 146, 60)

tb = slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(9), Inches(0.4))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "NOTRE APPROCHE : KAIZEN - Petites victoires = Grands resultats"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = RGBColor(96, 165, 250)
p.alignment = PP_ALIGN.CENTER

tb = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.3))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "3/18"
p.font.size = Pt(10)
p.font.color.rgb = RGBColor(147, 197, 253)
p.alignment = PP_ALIGN.CENTER

# ===== SLIDE 4 : VA / NVA =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(15, 45, 74)

tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.5))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "PARTIE 1 - VALEUR AJOUTEE"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)

tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(0.4))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "Ce que le client paie - et ce qu'il ne paie pas"
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(147, 197, 253)

# 3 categories
categories = [
    ("VA - 30%", "Valeur Ajoutee - Ce que le client PAIE", RGBColor(34, 197, 94)),
    ("NVA - 25%", "Non-VA Necessaire - Obligatoire (securite, admin)", RGBColor(251, 146, 60)),
    ("GASPILLAGE - 45%", "NVA Supprimable - Pur gaspillage a ELIMINER", RGBColor(239, 68, 68)),
]

y = 1.7
for pct, desc, color in categories:
    box = slide.shapes.add_shape(1, Inches(0.8), Inches(y), Inches(2.5), Inches(0.7))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(3.5), Inches(y + 0.1), Inches(5.5), Inches(0.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(255, 255, 255)

    y += 1.0

tb = slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(9), Inches(0.5))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "REDUIRE LE GASPILLAGE DE 45% A 20% = +2h/personne/jour = 25% DE CAPACITE SUPPLEMENTAIRE"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = RGBColor(96, 165, 250)
p.alignment = PP_ALIGN.CENTER

tb = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.3))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "4/18"
p.font.size = Pt(10)
p.font.color.rgb = RGBColor(147, 197, 253)
p.alignment = PP_ALIGN.CENTER

# ===== SLIDE 5 : 7 MUDA =====
add_simple_slide(prs, "PARTIE 1 - LES 7 MUDA",
    ["1. SURPRODUCTION : Rapports non lus, Stocks pieces excessifs",
     "2. ATTENTE : Technicien attend piece, Attente autorisation",
     "3. TRANSPORT : Pieces stock loin, Allers-retours magasin",
     "4. TRAITEMENT EXCESSIF : Rapports trop detailles, 3 validations",
     "5. STOCKS : Pieces qui dorment, Obsolescence",
     "6. MOUVEMENTS : Chercher outillage, Deplacements inutiles",
     "7. DEFAUTS : Reprise, Retouche, Mauvais diagnostic",
     "+ 8eme MUDA : TALENT IGNORE - Ne pas ecouter ceux qui font le travail"],
    "MUDA = Gaspillage - Les 7 gaspillages a eliminer en maintenance")

# ===== PARTIE 2 : DEMARCHES =====
# SLIDE 6 : PDCA
add_simple_slide(prs, "PARTIE 2 - DEMARCHES : PDCA",
    ["P - PLAN : Definir l'objectif, Analyser les causes, Planifier l'action",
     "D - DO : Tester la solution, Perimetre limite, Mesurer en continu",
     "C - CHECK : Mesurer les resultats, Analyser les ecarts, Valider l'efficacite",
     "A - ACT : Standardiser, Capitaliser, Etendre aux autres equipes"],
    "W. Edwards Deming 1950 - Cycle d'amelioration quotidienne")

# SLIDE 7 : DMAIC
add_simple_slide(prs, "PARTIE 2 - DEMARCHES : DMAIC",
    ["D - DEFINIR : Charte projet, QQOQCCP, Scope + Objectifs",
     "M - MESURER : Collecter donnees, Baseline, Indicateurs",
     "A - ANALYSER : 5 Pourquoi, Ishikawa, Donnees statistiques",
     "I - IMPROVE : Generer solutions, Pilotes, Test",
     "C - CONTROL : Standardiser, Suivi KPI, Documentation"],
    "Six Sigma - Motorola 1986 - Projets complexes 1-3 mois")

# SLIDE 8 : 8D
add_simple_slide(prs, "PARTIE 2 - DEMARCHES : 8D",
    ["D1 - EQUIPE : Multidisciplinaire (Meca + Elec + Prod)",
     "D2 - DEFINIR : Description precise, QQOQCCP complet",
     "D3 - ACTION IMMEDIATE : Stopper l'hemorragie, Mesure conservatoire",
     "D4 - CAUSES RACINES : 5 Pourquoi + Ishikawa, Preuves",
     "D5 - SOLUTIONS : Impact vs Faisabilite, Permanent",
     "D6 - IMPLEMENTER : Plan + Responsables, Delais",
     "D7 - PREVENTIVES : Empêcher récurrence, Autres equipements",
     "D8 - CAPITALISER : Partager leçons, Base de connaissances"],
    "Ford Motor Company 1987 - Incident grave - 3 jours a 2 semaines")

# ===== PARTIE 3 : OUTILS TERRAIN =====
# SLIDE 9 : QQOQCCP
add_simple_slide(prs, "PARTIE 3 - OUTILS TERRAIN : QQOQCCP",
    ["Q - QUI : Equipe A, Shift matin, Technicien: M. Benali",
     "Q - QUOI : Vibration anormale broyeur, Bruit inhabituel",
     "O - OU : Broyeur BK3, Palier 2, Cote moteur",
     "Q - QUAND : Demarrage production, 6h du matin, Depuis 3 jours",
     "C - COMMENT : Vibration 12mm/s, Alarme rouge, Temperature +15C",
     "C - COMBIEN : 3x/semaine, -2h/arret, -500t production",
     "P - POURQUOI : Perte production 500t/an, Risque casseur, Cout 20k Dh"],
    "Methode journalistique - 30 minutes - Pour bien definir un probleme")

# SLIDE 10 : 5 POURQUOI
add_simple_slide(prs, "PARTIE 3 - OUTILS TERRAIN : 5 POURQUOI",
    ["PROBLEME : Moteur convoyeur L4 en panne -> Arret ligne 4h, Perte 40t, Cout 8k Dh",
     "1. Pourquoi panne ? -> Surchauffe moteur",
     "2. Pourquoi surchauffe ? -> Ventilation encrassee",
     "3. Pourquoi encrassee ? -> Filtre absent depuis 3 mois",
     "4. Pourquoi absent ? -> Pas de plan maintenance preventive",
     "5. CAUSE RACINE -> Systeme de gestion MP inexistant",
     "SOLUTION : Implémenter plan MP + GEMBA mensuel + Check-liste",
     "RESULTAT : 0 panne ventilation en 6 mois, Gain 24h/an, Economie 12k Dh"],
    "Taiichi Ohno - Toyota 1950 - 15 minutes - Trouver la cause racine")

# SLIDE 11 : SMED
add_simple_slide(prs, "PARTIE 3 - OUTILS TERRAIN : SMED",
    ["PRINCIPE : Transformer taches INTERNES (machine arrêtée) en EXTERNES (machine en marche)",
     "INTERNE : Démontage, Nettoyage, Installation, Alignement, Test",
     "EXTERNE : Préparer outils, Préchauffer, À portée main, Check-list prête",
     "EXEMPLE OCP - Changement mâles concasseur :",
     "  AVANT : 4h 30min (INTERNE)",
     "  APRES : 2h 15min (SMED)",
     "  GAIN : -50%, 2h/fois, 150h/an sur 3 concasseurs",
     "ACTIONS : Chariot outils dédié + Check-list + Préparation shift avant"],
    "Single Minute Exchange of Die - Shigeo Shingo 1950 - -50% temps changement")

# SLIDE 12 : BRAINSTORMING
add_simple_slide(prs, "PARTIE 3 - OUTILS TERRAIN : BRAINSTORMING",
    ["CLASSIQUE : 6-10 pers, 20-40 min - Réunion équipe, Animateur, Tableau blanc",
     "STAND-UP : 5-8 pers, 10-15 min - Debout, Rapide, 1 problème unique",
     "BRAIN WRITING : 6 pers, 30 min - Chacun écrit 3 idées, Passe feuille",
     "SOLO : 1 pers, 15 min - Seul, Note TOUS les gaspillages vus",
     "EXEMPLE SOLO : Technicien note 12 gaspillages pendant son shift -> Detection précoce 3 pannes évitées"],
    "4 formats adaptés terrain maintenance")

# SLIDE 13 : 5G + ANDON
add_simple_slide(prs, "PARTIE 3 - OUTILS 2ème GENERATION : 5G + ANDON",
    ["5G - Les 5 Gemba :",
     "  1. GEMBA - Aller sur terrain, Atelier, Ligne de production",
     "  2. GEMBUTSU - Observer l'équipement réel, Toucher, Ecouter",
     "  3. GENJITSU - Données réelles, MES, Pas de supposition",
     "  4. GENRI - Principes ingénierie, Physique, Mécanique",
     "  5. GEMBA WALK - 15 min/jour, Noter 3 MUDA",
     "ANDON - Management Visuel :",
     "  VERT - Equipement OK, Production normale",
     "  ORANGE - Vigilance, Planifié maintenance",
     "  ROUGE - Panne, Action immédiate requise",
     "EXEMPLE OCP : Tableau ANDON salle contrôle - 20 équipements - Temps réaction panne -40%"],
    "Management visuel - 15 min/jour")

# SLIDE 14 : 5G VS GENGSOKU
add_simple_slide(prs, "PARTIE 3 - COMPARAISON 5G / GENGSOKU",
    ["5G - OBSERVER : Aller sur terrain, Voir l'équipement, Données réelles, Principes, Walk 15 min/jour",
     "  But : OBSERVER la réalité, Comprendre ce qui se passe réellement sur le terrain",
     "GENGSOKU - EXPÉRIMENTER : Faire soi-même, Essayer l'opération, Ressentir la difficulté, Identifier freins réels",
     "  But : VIVRE le problème, Expérimenter ce que les opérateurs subissent",
     "5G + GENGSOKU = Complémentaires pour une vraie compréhension",
     "EXEMPLE : 5G -> Je vois le technicien chercher l'outil",
     "  GENGSOKU -> Je cherche l'outil moi-même et je perds 15 min",
     "  ACTION : Chariot outillage mobile"],
    "Comprendre l'écart pour ajuster votre pratique")

# ===== PARTIE 4 : GAINS + PLAN =====
# SLIDE 15 : GAINS
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(15, 45, 74)

tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.5))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "PARTIE 4 - GAINS"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)

tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(0.4))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "Ce que vous gagnez réellement - Résultats mesurés OCP + industrie"
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(147, 197, 253)

gains = [
    ("-40%", "Pannes récurrentes", "En 6 mois avec PDCA + 5P"),
    ("+2h/j", "Disponibilité", "Par équipement, Sans investissement"),
    ("-50%", "Temps d'intervention", "Avec SMED + Organisation"),
    ("0 Dh", "D'investissement", "Juste votre regard, 15 min/jour"),
]

positions = [(0.6, 1.6), (5.0, 1.6), (0.6, 3.4), (5.0, 3.4)]
colors = [RGBColor(34, 197, 94), RGBColor(59, 130, 246), RGBColor(168, 85, 247), RGBColor(251, 146, 60)]

for i, (value, title, desc) in enumerate(gains):
    x, y = positions[i]

    box = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(4.2), Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(22, 43, 66)
    box.line.color.rgb = colors[i]
    box.line.width = Pt(2)

    tb = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.2), Inches(3.8), Inches(0.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = value
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = colors[i]
    p.alignment = PP_ALIGN.CENTER

    tb = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.7), Inches(3.8), Inches(0.35))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    tb = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 1.05), Inches(3.8), Inches(0.35))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER

timeline = [
    "1ere SEMAINE : Identifier MUDA, Gemba Walk",
    "1er MOIS : Premiers resultats visibles",
    "3 MOIS : Equipe autonome",
    "6 MOIS : -40% pannes recurrentes",
]

y = 5.2
for item in timeline:
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(8.4), Inches(0.25))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = item
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(255, 255, 255)
    y += 0.3

tb = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.3))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "15/18"
p.font.size = Pt(10)
p.font.color.rgb = RGBColor(147, 197, 253)
p.alignment = PP_ALIGN.CENTER

# SLIDE 16 : PLAN D'ACTION 18-22 MAI
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(15, 45, 74)

tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.5))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "PARTIE 4 - PLAN D'ACTION"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)

tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(0.4))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "Semaine pratique - Du 18 au 22 Mai 2026"
p.font.size = Pt(18)
p.font.color.rgb = RGBColor(147, 197, 253)

week_plan = [
    ("LUNDI 18", "5G - Gemba Walk", "15 min terrain, Noter 3 MUDA observés", "3 MUDA identifiés", RGBColor(59, 130, 246)),
    ("MARDI 19", "QQOQCCP + 5 Pourquoi", "30 min sur MUDA critique, Équipe en stand-up", "Cause racine trouvée", RGBColor(34, 197, 94)),
    ("MERCREDI 20", "Brainstorming", "10 min stand-up, Équipe propose solutions", "Solution choisie", RGBColor(168, 85, 247)),
    ("JEUDI 21", "PDCA - Do", "Tester solution, Mesurer avant/après", "Résultat mesuré", RGBColor(251, 146, 60)),
    ("VENDREDI 22", "Standardisation", "Si succès : Standard, Partager aux autres équipes", "Standard diffusé", RGBColor(34, 197, 94)),
]

y = 1.6
for day, tool, action, result, color in week_plan:
    box = slide.shapes.add_shape(1, Inches(0.6), Inches(y), Inches(8.8), Inches(0.9))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(22, 43, 66)
    box.line.color.rgb = color

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(y + 0.1), Inches(1.5), Inches(0.3))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = day
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = color

    tb = slide.shapes.add_textbox(Inches(2.4), Inches(y + 0.1), Inches(2.0), Inches(0.3))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = tool
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(y + 0.45), Inches(4.5), Inches(0.35))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = action
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(200, 200, 200)

    box2 = slide.shapes.add_shape(1, Inches(5.5), Inches(y + 0.15), Inches(3.7), Inches(0.6))
    box2.fill.solid()
    box2.fill.fore_color.rgb = color
    box2.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(5.7), Inches(y + 0.25), Inches(3.3), Inches(0.4))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = f"OK {result}"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    y += 1.1

tb = slide.shapes.add_textbox(Inches(0.6), Inches(7.2), Inches(8.8), Inches(0.4))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "OBJECTIF SEMAINE : 1 KAIZEN testé par équipe = 5 améliorations concrètes OCP"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = RGBColor(96, 165, 250)
p.alignment = PP_ALIGN.CENTER

tb = slide.shapes.add_textbox(Inches(0.5), Inches(7.6), Inches(9), Inches(0.3))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "16/18"
p.font.size = Pt(10)
p.font.color.rgb = RGBColor(147, 197, 253)
p.alignment = PP_ALIGN.CENTER

# SLIDE 17 : MÉMO OUTILS
add_simple_slide(prs, "ANNEXE - MÉMO OUTILS TERRAIN",
    ["PDCA : Amelioration quotidienne, Cycle 1 semaine",
     "QQOQCCP : Definir probleme, 30 min",
     "5 Pourquoi : Cause racine, 15 min",
     "SMED : Reduire temps changement, Projet 1 mois",
     "Brainstorming : Generer solutions, 10-40 min",
     "5G/Andon : Management visuel, 15 min/jour",
     "8D : Incident grave, 3j-2sem",
     "DMAIC : Projet complexe, 1-3 mois",
     "CHECKLIST QUOTIDIEN : Gemba Walk 15min, Verifier ANDON, Briefing équipe, Valider 1 KAIZEN"],
    "A emporter - Guide rapide terrain")

# SLIDE 18 : CONCLUSION
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(15, 45, 74)

tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(0.8))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "\"Lean Maintenance : Un voyage, pas une destination.\""
p.font.size = Pt(24)
p.font.italic = True
p.font.color.rgb = RGBColor(96, 165, 250)
p.alignment = PP_ALIGN.CENTER

principes = [
    ("Commencez PETIT", "1 MUDA a la fois"),
    ("Commencez MAINTENANT", "Lundi 18 mai"),
    ("Capitalisez TOUJOURS", "Partagez les succès"),
]

y = 2.6
for principe, desc in principes:
    box = slide.shapes.add_shape(1, Inches(1.5), Inches(y), Inches(7.0), Inches(0.6))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(22, 43, 66)
    box.line.color.rgb = RGBColor(96, 165, 250)

    tb = slide.shapes.add_textbox(Inches(1.7), Inches(y + 0.1), Inches(3.0), Inches(0.4))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = principe
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(96, 165, 250)

    tb = slide.shapes.add_textbox(Inches(4.8), Inches(y + 0.1), Inches(3.2), Inches(0.4))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(255, 255, 255)

    y += 0.8

tb = slide.shapes.add_textbox(Inches(0.6), Inches(5.2), Inches(8.8), Inches(0.5))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "Votre equipe vous attend. Les outils sont prets. La premiere amelioration commence lundi."
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = RGBColor(34, 197, 94)
p.alignment = PP_ALIGN.CENTER

tb = slide.shapes.add_textbox(Inches(2.0), Inches(6.0), Inches(6.0), Inches(0.4))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "Direction Maintenance OCP - Formation KAIZEN"
p.font.size = Pt(12)
p.font.color.rgb = RGBColor(200, 200, 200)
p.alignment = PP_ALIGN.CENTER

tb = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.3))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "18/18"
p.font.size = Pt(10)
p.font.color.rgb = RGBColor(147, 197, 253)
p.alignment = PP_ALIGN.CENTER

# Sauvegarder
output_path = "C:/Users/HP/Desktop/formation/Formation_Lean_KAIZEN_OCP_V2.pptx"
prs.save(output_path)
print(f"Presentation creee : {output_path}")
