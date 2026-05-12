# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Couleurs du theme Lean
BLEU_FONCE = RGBColor(0, 56, 102)
BLEU_PRINCIPAL = RGBColor(0, 75, 135)
BLEU_CLAIR = RGBColor(96, 165, 250)
BLEU_PALE = RGBColor(147, 197, 253)
VERT_FONCE = RGBColor(22, 101, 52)
VERT = RGBColor(41, 139, 58)
VERT_CLAIR = RGBColor(110, 185, 46)
BLANC = RGBColor(255, 255, 255)
GRIS_SOMBRE = RGBColor(17, 24, 39)
GRIS_MOYEN = RGBColor(55, 65, 81)
JAUNE = RGBColor(217, 119, 6)
ORANGE = RGBColor(237, 125, 49)
ROUGE = RGBColor(185, 28, 28)

def create_slide_1_brain_vs_chaos():
    """Modele 1: Le chaos cache - Contraste cerveau vs desorganisation"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Fond bleu
    bg = slide.shapes.add_shape(1, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = BLEU_FONCE
    bg.line.fill.background()

    # Titre
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7))
    tf = title.text_frame
    tf.text = "CE QUE VOUS VOYEZ vs CE QUI SE PASSE VRAIMENT"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = BLANC
    p.alignment = PP_ALIGN.CENTER

    # Zone visuelle - Split
    # Gauche - La réalité du terrain
    left_bg = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(6), Inches(5.5))
    left_bg.fill.solid()
    left_bg.fill.fore_color.rgb = RGBColor(30, 35, 45)
    left_bg.line.color.rgb = ROUGE
    left_bg.line.width = Pt(2)

    left_title = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(5.6), Inches(0.5))
    tf = left_title.text_frame
    tf.text = "URGENCE DU JOUR"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ROUGE
    p.alignment = PP_ALIGN.CENTER

    problems = [
        ("Panne X encore une fois"),
        ("Recherche empirique de la solution"),
        ("Pièce manquante en stock"),
        ("Technicien qui part avec le savoir")
    ]

    y_pos = 2.2
    for prob in problems:
        box = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(5), Inches(0.5))
        tf = box.text_frame
        tf.text = f"• {prob}"
        p = tf.paragraphs[0]
        p.font.size = Pt(12)
        p.font.color.rgb = BLANC
        y_pos += 0.65

    # Message impact
    impact = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(5), Inches(1.5))
    tf = impact.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "On dit: 'C'est le vieil equipement'"
    p.font.size = Pt(13)
    p.font.color.rgb = ORANGE
    p.italic = True
    p2 = tf.add_paragraph()
    p2.text = "\nRealite: On n'a jamais analyse le vrai probleme."
    p2.font.size = Pt(11)
    p2.font.color.rgb = BLEU_PALE

    # Droite - Ce que le Lean voit
    right_bg = slide.shapes.add_shape(1, Inches(6.8), Inches(1.3), Inches(6), Inches(5.5))
    right_bg.fill.solid()
    right_bg.fill.fore_color.rgb = RGBColor(35, 45, 35)
    right_bg.line.color.rgb = VERT_CLAIR
    right_bg.line.width = Pt(2)

    right_title = slide.shapes.add_textbox(Inches(7), Inches(1.5), Inches(5.6), Inches(0.5))
    tf = right_title.text_frame
    tf.text = "CE QUE LE LEAN VOIT"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = VERT_CLAIR
    p.alignment = PP_ALIGN.CENTER

    solutions = [
        ("Cause racine jamais identifiee"),
        ("Solution standardisee existe"),
        ("Stock visuel gere"),
        ("Savoir documente et partage")]
    y_pos = 2.2
    for sol in solutions:
        box = slide.shapes.add_textbox(Inches(7.3), Inches(y_pos), Inches(5), Inches(0.5))
        tf = box.text_frame
        tf.text = f"✓ {sol}"
        p = tf.paragraphs[0]
        p.font.size = Pt(12)
        p.font.color.rgb = BLANC
        y_pos += 0.65

    # Message espoir
    espoir = slide.shapes.add_textbox(Inches(7.3), Inches(5), Inches(5), Inches(1.5))
    tf = espoir.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "On dit: 'On ne savait pas'"
    p.font.size = Pt(13)
    p.font.color.rgb = VERT_CLAIR
    p.italic = True
    p2 = tf.add_paragraph()
    p2.text = "\nRealite: On n'avait jamais regarde."
    p2.font.size = Pt(11)
    p2.font.color.rgb = BLEU_PALE

    # Footer
    footer = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(12.333), Inches(0.3))
    tf = footer.text_frame
    tf.text = "Le Lean ne change pas vos equipements. Il change votre regard."
    p = tf.paragraphs[0]
    p.font.size = Pt(12)
    p.font.color.rgb = BLEU_PALE
    p.alignment = PP_ALIGN.CENTER

    prs.save('Modele1_Chaos_Cache.pptx')
    print("OK - Modele 1 cree: Chaos Cache - Contraste cerveau vs desorganisation")

def create_slide_2_triangle_risk():
    """Modele 2: Le triangle du risque - Approche visuelle danger"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Fond sombre dramatique
    bg = slide.shapes.add_shape(1, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(20, 20, 25)
    bg.line.fill.background()

    # Titre impactant
    title = slide.shapes.add_textbox(Inches(1), Inches(0.4), Inches(11.333), Inches(0.8))
    tf = title.text_frame
    tf.text = "3 PIEGES QUI COUTENT CHAQUE JOUR"
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.alignment = PP_ALIGN.CENTER

    # 3 pieges
    pieges = [
        ("1", "L'URGENCE", "Devenue normalite", ROUGE,
         "Chaque jour commence par: QUOI DE CASSE?"),
        ("2", "LE REEMPLOI", "Du meme erreur", JAUNE,
         "Meme panne. Meme cause. Meme solution temporaire."),
        ("3", "LA DEPENDANCE", "Aux heros de service", VERT_FONCE,
         "Si un expert part, 20 ans de savoir partent.")
    ]

    y_pos = 1.6
    for num, titre, soustitre, couleur, desc in pieges:
        # Box
        box = slide.shapes.add_shape(1, Inches(1), Inches(y_pos), Inches(11.333), Inches(1.1))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(35, 35, 40)
        box.line.color.rgb = couleur
        box.line.width = Pt(3)

        # Numero
        num_box = slide.shapes.add_textbox(Inches(1.3), Inches(y_pos + 0.2), Inches(0.6), Inches(0.7))
        tf = num_box.text_frame
        tf.text = num
        p = tf.paragraphs[0]
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = couleur

        # Titre
        titre_box = slide.shapes.add_textbox(Inches(2.1), Inches(y_pos + 0.15), Inches(3), Inches(0.35))
        tf = titre_box.text_frame
        tf.text = titre
        p = tf.paragraphs[0]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = BLANC

        # Sous-titre
        st_box = slide.shapes.add_textbox(Inches(2.1), Inches(y_pos + 0.55), Inches(3), Inches(0.35))
        tf = st_box.text_frame
        tf.text = soustitre
        p = tf.paragraphs[0]
        p.font.size = Pt(12)
        p.font.color.rgb = GRIS_MOYEN

        # Description
        desc_box = slide.shapes.add_textbox(Inches(5.5), Inches(y_pos + 0.2), Inches(6.5), Inches(0.7))
        tf = desc_box.text_frame
        tf.word_wrap = True
        tf.text = desc
        p = tf.paragraphs[0]
        p.font.size = Pt(12)
        p.font.color.rgb = BLEU_PALE

        y_pos += 1.25

    # Message d'espoir
    hope_bg = slide.shapes.add_shape(1, Inches(1), Inches(5.6), Inches(11.333), Inches(1.5))
    hope_bg.fill.solid()
    hope_bg.fill.fore_color.rgb = VERT_FONCE
    hope_bg.line.fill.background()

    hope_text = slide.shapes.add_textbox(Inches(1), Inches(5.9), Inches(11.333), Inches(1))
    tf = hope_text.text_frame
    tf.text = "LE LEAN CASSE CE TRIANGLE"
    p = tf.paragraphs[0]
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = BLANC
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "Pas besoin de budget. Juste un nouveau regard."
    p2.font.size = Pt(14)
    p2.font.color.rgb = VERT_CLAIR
    p2.alignment = PP_ALIGN.CENTER

    prs.save('Modele2_Triangle_Risque.pptx')
    print("OK - Modele 2 cree: Triangle du Risque - Approche visuelle danger")

def create_slide_3_brain_cost():
    """Modele 3: Cerveau solage - Message qualitatif impactant"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Fond bleu
    bg = slide.shapes.add_shape(1, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = BLEU_FONCE
    bg.line.fill.background()

    # Titre
    title = slide.shapes.add_textbox(Inches(2), Inches(0.4), Inches(9.333), Inches(0.8))
    tf = title.text_frame
    tf.text = "POURQUOI PERSONNE NE VOIT CE GASPILLAGE?"
    p = tf.paragraphs[0]
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = BLANC
    p.alignment = PP_ALIGN.CENTER

    # Zone centrale
    center_bg = slide.shapes.add_shape(1, Inches(1.5), Inches(1.6), Inches(10.333), Inches(2.5))
    center_bg.fill.solid()
    center_bg.fill.fore_color.rgb = RGBColor(40, 50, 70)
    center_bg.line.color.rgb = BLEU_CLAIR
    center_bg.line.width = Pt(3)

    # Message principal
    main_text = slide.shapes.add_textbox(Inches(2), Inches(2.2), Inches(9.333), Inches(0.6))
    tf = main_text.text_frame
    tf.word_wrap = True
    tf.text = "Parce qu'on a normalise l'anormal."
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = VERT_CLAIR
    p.alignment = PP_ALIGN.CENTER

    # Citation
    quote = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(9.333), Inches(0.8))
    tf = quote.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = '"On dirait que c\'est dans la nature de cette machine de tomber en panne."'
    p.font.size = Pt(14)
    p.font.color.rgb = BLANC
    p.alignment = PP_ALIGN.CENTER
    p.italic = True
    p2 = tf.add_paragraph()
    p2.text = "Cette phrase dit que votre equipe a abandonne."
    p2.font.size = Pt(12)
    p2.font.color.rgb = ORANGE
    p2.alignment = PP_ALIGN.CENTER

    # 3 verites
    verites = [
        ("L'ILLUSION", "On croit que c'est un probleme d'equipement. C'est un probleme de processus."),
        ("LE SILENCE", "Chaque technician sait par experience. Mais personne ne documente."),
        ("LE COUT", "Pas celui de la panne. Celui de ne JAMAIS apprendre de la panne.")
    ]

    y_pos = 4.4
    for titre, desc in verites:
        box = slide.shapes.add_shape(1, Inches(1), Inches(y_pos), Inches(11.333), Inches(0.7))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0, 85, 150)
        box.line.fill.background()

        box_text = slide.shapes.add_textbox(Inches(1.5), Inches(y_pos + 0.15), Inches(10.333), Inches(0.5))
        tf = box_text.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"• {titre}: {desc}"
        p.font.size = Pt(12)
        p.font.color.rgb = BLANC
        y_pos += 0.8

    # Footer
    footer = slide.shapes.add_textbox(Inches(1), Inches(6.9), Inches(11.333), Inches(0.4))
    tf = footer.text_frame
    tf.text = "Le Lean? C'est juste une autre facon de regarder ce qui existe deja."
    p = tf.paragraphs[0]
    p.font.size = Pt(12)
    p.font.color.rgb = BLEU_PALE
    p.alignment = PP_ALIGN.CENTER

    prs.save('Modele3_Cerveau_Impact.pptx')
    print("OK - Modele 3 cree: Cerveau Solage - Message qualitatif impactant")

def create_slide_4_before_after():
    """Modele 4: Avant/Apres dramatique"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Titre
    title = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(11.333), Inches(0.7))
    tf = title.text_frame
    tf.text = "UNE JOURNEE SANS LEAN = UNE JOURNEE EN FEU"
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = GRIS_SOMBRE
    p.alignment = PP_ALIGN.CENTER

    # Separator ligne
    line = slide.shapes.add_shape(1, Inches(6.4), Inches(1.2), Inches(0.1), Inches(5.2))
    line.fill.solid()
    line.fill.fore_color.rgb = GRIS_MOYEN
    line.line.fill.background()

    # Fleche centrale
    arrow_box = slide.shapes.add_textbox(Inches(5.8), Inches(3.6), Inches(1.8), Inches(0.5))
    tf = arrow_box.text_frame
    tf.text = "LEAN"
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = VERT
    p.alignment = PP_ALIGN.CENTER

    # GAUCHE - AVANT
    before_bg = slide.shapes.add_shape(1, 0, Inches(1.2), Inches(6.4), Inches(5.2))
    before_bg.fill.solid()
    before_bg.fill.fore_color.rgb = RGBColor(254, 226, 226)
    before_bg.line.fill.background()

    before_title = slide.shapes.add_textbox(Inches(0.3), Inches(1.4), Inches(5.8), Inches(0.4))
    tf = before_title.text_frame
    tf.text = "SANS LEAN - VOTRE ACTUEL"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ROUGE

    before_items = [
        "07:00 - Le planning est deja depasse",
        "09:30 - Panne non-documentee. Recherche empirique.",
        "11:00 - Piece manquante. 'C'est arrive la semaine derniere'",
        "14:00 - Meme equipement. Meme panne. 'Vieille machine'",
        "16:30 - Rapport? On a pas eu le temps",
        "17:00 - Demain sera pareil. 'C'est comme ca ici'"
    ]

    y_pos = 2
    for item in before_items:
        item_box = slide.shapes.add_textbox(Inches(0.4), Inches(y_pos), Inches(5.6), Inches(0.42))
        tf = item_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(10)
        p.font.color.rgb = GRIS_SOMBRE
        y_pos += 0.52

    # DROITE - APRES
    after_bg = slide.shapes.add_shape(1, Inches(6.5), Inches(1.2), Inches(6.833), Inches(5.2))
    after_bg.fill.solid()
    after_bg.fill.fore_color.rgb = RGBColor(220, 252, 231)
    after_bg.line.fill.background()

    after_title = slide.shapes.add_textbox(Inches(7), Inches(1.4), Inches(5.8), Inches(0.4))
    tf = after_title.text_frame
    tf.text = "AVEC LEAN - LE POSSIBLE"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = VERT_FONCE

    after_items = [
        "07:00 - Brief 5min. Aujourd'hui, on suit le plan.",
        "09:30 - Historique consulte. Solution documentee en 3 min.",
        "11:00 - Stock gere visuellement. La piece etait la.",
        "14:00 - Analyse racine. 'Cette panne ne reviendra plus'",
        "16:30 - 5min de capture d'apprentissage. Capitalisation.",
        "17:00 - Demain sera meilleur. 'On a progresse'"
    ]

    y_pos = 2
    for item in after_items:
        item_box = slide.shapes.add_textbox(Inches(7.2), Inches(y_pos), Inches(5.8), Inches(0.42))
        tf = item_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(10)
        p.font.color.rgb = GRIS_SOMBRE
        y_pos += 0.52

    # Footer
    footer = slide.shapes.add_shape(1, 0, Inches(6.5), Inches(13.333), Inches(1))
    footer.fill.solid()
    footer.fill.fore_color.rgb = VERT
    footer.line.fill.background()

    footer_text = slide.shapes.add_textbox(Inches(1), Inches(6.8), Inches(11.333), Inches(0.4))
    tf = footer_text.text_frame
    tf.text = "La difference n'est pas dans les moyens. Elle est dans la methode."
    p = tf.paragraphs[0]
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = BLANC
    p.alignment = PP_ALIGN.CENTER

    prs.save('Modele4_Before_After.pptx')
    print("OK - Modele 4 cree: Avant/Apres Dramatique - Comparaison visuelle")

def create_slide_5_question_puissante():
    """Modele 5: Question puissante qui fait reflechir"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Fond noir profond
    bg = slide.shapes.add_shape(1, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(10, 10, 15)
    bg.line.fill.background()

    # Bordure accent
    border_top = slide.shapes.add_shape(1, 0, 0, Inches(13.333), Inches(0.1))
    border_top.fill.solid()
    border_top.fill.fore_color.rgb = VERT_CLAIR
    border_top.line.fill.background()

    border_bottom = slide.shapes.add_shape(1, 0, Inches(7.4), Inches(13.333), Inches(0.1))
    border_bottom.fill.solid()
    border_bottom.fill.fore_color.rgb = VERT_CLAIR
    border_bottom.line.fill.background()

    # Question principale
    question = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.333), Inches(1.5))
    tf = question.text_frame
    tf.word_wrap = True
    tf.text = "Si vous connaissiez le cout reel"
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = BLANC
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "de 'faire comme d'habitude'..."
    p2.font.size = Pt(32)
    p2.font.bold = True
    p2.font.color.rgb = ORANGE
    p2.alignment = PP_ALIGN.CENTER

    # Pause visuelle
    pause = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(11.333), Inches(0.5))
    tf = pause.text_frame
    tf.text = "—"
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.color.rgb = GRIS_MOYEN
    p.alignment = PP_ALIGN.CENTER

    # Deuxieme partie
    question2 = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11.333), Inches(1))
    tf = question2.text_frame
    tf.word_wrap = True
    tf.text = "changeriez-vous quelque chose?"
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = VERT_CLAIR
    p.alignment = PP_ALIGN.CENTER

    # Zone de reponse
    reponse_bg = slide.shapes.add_shape(1, Inches(3), Inches(5.1), Inches(7.333), Inches(1.4))
    reponse_bg.fill.solid()
    reponse_bg.fill.fore_color.rgb = BLEU_FONCE
    reponse_bg.line.color.rgb = BLEU_CLAIR
    reponse_bg.line.width = Pt(2)

    reponse = slide.shapes.add_textbox(Inches(3.5), Inches(5.3), Inches(6.333), Inches(1))
    tf = reponse.text_frame
    tf.word_wrap = True
    tf.text = "Le probleme n'est pas que vous ne savez PAS.\n\nLe probleme est que vous ne CALCULEZ PAS."
    p = tf.paragraphs[0]
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = BLANC
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.font.size = Pt(15)
    p2.font.color.rgb = BLEU_PALE
    p2.alignment = PP_ALIGN.CENTER

    # Appel a l'action
    cta = slide.shapes.add_textbox(Inches(1), Inches(6.7), Inches(11.333), Inches(0.5))
    tf = cta.text_frame
    tf.text = "Les 10 prochaines minutes vont changer votre regard sur vos ateliers."
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.color.rgb = GRIS_MOYEN
    p.alignment = PP_ALIGN.CENTER

    prs.save('Modele5_Question_Impact.pptx')
    print("OK - Modele 5 cree: Question Puissante - Message reflexif")

# Creer les 5 modeles
if __name__ == "__main__":
    create_slide_1_brain_vs_chaos()
    create_slide_2_triangle_risk()
    create_slide_3_brain_cost()
    create_slide_4_before_after()
    create_slide_5_question_puissante()

    print("\nSUCCESS - Tous les 5 modeles ont ete crees!")
